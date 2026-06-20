"""Deterministic outline -> PPTX renderer (epic #218, child #221).

Kualitas visual berasal dari renderer ini (design tokens + layout terkontrol),
bukan dari model AI. Renderer menghasilkan deck 16:9 dengan logo, header, footer
personalisasi, dan 5 template visual resmi. Tidak ada akses internet.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from datetime import datetime
from io import BytesIO
from typing import Any, Iterable, Mapping

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.services.presentation_assets import (
    BRAND_NAME,
    add_brand_logo,
    add_icon,
)

# ── Batas/validasi ────────────────────────────────────────────────────────
SLIDE_COUNT_MIN = 3
SLIDE_COUNT_MAX = 20
MAX_CONTENT_SLIDES = 18
MAX_BULLETS_PER_SLIDE = 7
BULLET_MAX_CHARS = 240
TITLE_MAX_CHARS = 160
SLIDE_TITLE_MAX_CHARS = 120

# Slide 16:9 (widescreen).
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.62)


@dataclass(frozen=True)
class Palette:
    primary: str
    accent: str
    background: str
    surface: str
    text: str
    muted: str
    on_primary: str = "FFFFFF"


@dataclass(frozen=True)
class PresentationTemplate:
    key: str
    label: str
    palette: Palette
    heading_font: str
    body_font: str
    cover_style: str  # band | minimal | split | sidebar | frame


PRESENTATION_TEMPLATES: dict[str, PresentationTemplate] = {
    "resmi_klasik": PresentationTemplate(
        key="resmi_klasik",
        label="Resmi Klasik",
        palette=Palette(
            primary="1F3A5F", accent="C9A227", background="FFFFFF",
            surface="F4F6FB", text="1A2230", muted="5B6B82",
        ),
        heading_font="Georgia",
        body_font="Arial",
        cover_style="band",
    ),
    "modern_minimal": PresentationTemplate(
        key="modern_minimal",
        label="Modern Minimal",
        palette=Palette(
            primary="111827", accent="0EA5A4", background="FBFBFD",
            surface="F1F5F9", text="0F172A", muted="64748B",
        ),
        heading_font="Calibri",
        body_font="Calibri",
        cover_style="minimal",
    ),
    "executive_brief": PresentationTemplate(
        key="executive_brief",
        label="Executive Brief",
        palette=Palette(
            primary="5B1A22", accent="E0A106", background="FBF7F0",
            surface="F3EADB", text="241015", muted="7A5A52",
        ),
        heading_font="Georgia",
        body_font="Arial",
        cover_style="split",
    ),
    "data_tabel": PresentationTemplate(
        key="data_tabel",
        label="Data & Tabel",
        palette=Palette(
            primary="1E3A8A", accent="0891B2", background="FFFFFF",
            surface="EEF2FF", text="0F1B33", muted="55617A",
        ),
        heading_font="Calibri",
        body_font="Arial",
        cover_style="sidebar",
    ),
    "kegiatan_dokumentasi": PresentationTemplate(
        key="kegiatan_dokumentasi",
        label="Kegiatan & Dokumentasi",
        palette=Palette(
            primary="14532D", accent="84CC16", background="FBFDF7",
            surface="ECF5E2", text="14210F", muted="51664A",
        ),
        heading_font="Calibri",
        body_font="Arial",
        cover_style="frame",
    ),
}

# Ikon lokal per posisi slide (urut untuk variasi visual ringan).
_CONTENT_ICON_CYCLE = ("agenda", "summary", "key_points", "data", "activity")


@dataclass
class SlideContent:
    title: str
    bullets: list[str] = field(default_factory=list)
    layout: str = "content"


@dataclass
class PresentationRender:
    filename: str
    content: bytes
    template: str
    slide_count: int


def template_choices() -> list[dict[str, str]]:
    """Daftar template visual untuk UI/kontrak (label-friendly)."""
    return [{"key": t.key, "label": t.label} for t in PRESENTATION_TEMPLATES.values()]


def normalize_template(visual_template: str | None) -> PresentationTemplate:
    key = (visual_template or "").strip().lower().replace("-", "_").replace(" ", "_")
    template = PRESENTATION_TEMPLATES.get(key)
    if template is None:
        raise ValueError(
            "Template visual tidak dikenal. Pilihan: "
            + ", ".join(PRESENTATION_TEMPLATES.keys())
        )
    return template


def _clean_text(value: Any, *, limit: int) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip()
    if len(text) > limit:
        text = text[: limit - 1].rstrip() + "\u2026"
    return text


def _clean_multiline(value: Any, *, limit: int) -> str:
    text = re.sub(r"[ \t]+", " ", str(value or "")).strip()
    if len(text) > limit:
        text = text[: limit - 1].rstrip() + "\u2026"
    return text


def _coerce_bullets(raw_bullets: Any) -> list[str]:
    if raw_bullets is None:
        return []
    if isinstance(raw_bullets, str):
        candidates: Iterable[Any] = raw_bullets.splitlines()
    elif isinstance(raw_bullets, Iterable):
        candidates = raw_bullets
    else:
        candidates = [raw_bullets]

    bullets: list[str] = []
    for candidate in candidates:
        cleaned = _clean_text(candidate, limit=BULLET_MAX_CHARS)
        if cleaned:
            bullets.append(cleaned)
        if len(bullets) >= MAX_BULLETS_PER_SLIDE:
            break
    return bullets


def _normalize_outline(outline: Any) -> list[SlideContent]:
    slides: list[SlideContent] = []
    if not isinstance(outline, Iterable) or isinstance(outline, (str, bytes)):
        return slides

    for index, raw in enumerate(outline):
        if isinstance(raw, Mapping):
            title = _clean_text(
                raw.get("title") or raw.get("heading") or "",
                limit=SLIDE_TITLE_MAX_CHARS,
            )
            bullets = _coerce_bullets(
                raw.get("bullets") or raw.get("points") or raw.get("content")
            )
            layout = str(raw.get("layout") or "content").strip().lower()
        else:
            title = _clean_text(raw, limit=SLIDE_TITLE_MAX_CHARS)
            bullets = []
            layout = "content"

        if not title and not bullets:
            continue
        if not title:
            title = f"Bagian {index + 1}"

        slides.append(SlideContent(title=title, bullets=bullets, layout=layout))
        if len(slides) >= MAX_CONTENT_SLIDES:
            break

    return slides


def _clamp_slide_count(slide_count: Any) -> int | None:
    if slide_count is None:
        return None
    try:
        value = int(slide_count)
    except (TypeError, ValueError):
        return None
    return max(SLIDE_COUNT_MIN, min(SLIDE_COUNT_MAX, value))


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.lstrip("#").upper())


def _slugify(value: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", value.lower()).strip("-")
    return slug or "presentasi-ista"


def render_presentation(
    *,
    title: str,
    visual_template: str,
    outline: Any = None,
    subtitle: str | None = None,
    audience: str | None = None,
    header: str | None = None,
    footer: str | None = None,
    presenter: str | None = None,
    unit: str | None = None,
    slide_count: Any = None,
    date: str | None = None,
) -> PresentationRender:
    clean_title = _clean_text(title, limit=TITLE_MAX_CHARS)
    if not clean_title:
        raise ValueError("Judul presentasi wajib diisi.")

    template = normalize_template(visual_template)
    palette = template.palette

    content_slides = _normalize_outline(outline)
    target = _clamp_slide_count(slide_count)
    if target is not None:
        # cover + content + closing == target.
        max_content = max(1, target - 2)
        content_slides = content_slides[:max_content]

    clean_header = _clean_text(header or BRAND_NAME, limit=120)
    clean_footer = _clean_multiline(footer or BRAND_NAME, limit=160)
    clean_subtitle = _clean_text(subtitle or "", limit=200)
    clean_audience = _clean_text(audience or "", limit=160)
    clean_presenter = _clean_text(presenter or "", limit=160)
    clean_unit = _clean_text(unit or "", limit=160)
    clean_date = _clean_text(date or datetime.now().strftime("%d %B %Y"), limit=80)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _build_cover_slide(
        prs,
        template,
        title=clean_title,
        subtitle=clean_subtitle,
        audience=clean_audience,
        presenter=clean_presenter,
        unit=clean_unit,
        date=clean_date,
    )

    total_pages = len(content_slides) + 1  # +1 closing
    for index, slide_content in enumerate(content_slides):
        _build_content_slide(
            prs,
            template,
            slide_content,
            header=clean_header,
            footer=clean_footer,
            page_number=index + 1,
            total_pages=total_pages,
        )

    _build_closing_slide(
        prs,
        template,
        header=clean_header,
        footer=clean_footer,
        presenter=clean_presenter,
        unit=clean_unit,
        page_number=total_pages,
        total_pages=total_pages,
    )

    buffer = BytesIO()
    prs.save(buffer)

    return PresentationRender(
        filename=f"{_slugify(clean_title)}.pptx",
        content=buffer.getvalue(),
        template=template.key,
        slide_count=len(prs.slides._sldIdLst),
    )


# ── Primitives ──────────────────────────────────────────────────────────────
def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill_rect(slide, hex_color: str, left, top, width, height, *, name: str | None = None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_color)
    shape.line.fill.background()
    shape.shadow.inherit = False
    if name:
        shape.name = name
    return shape


def _background(slide, hex_color: str):
    _fill_rect(slide, hex_color, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, name="bg")


def _textbox(
    slide,
    text: str,
    *,
    left,
    top,
    width,
    height,
    size: int,
    color: str,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = "Arial",
    anchor=MSO_ANCHOR.TOP,
    line_spacing: float = 1.0,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = _rgb(color)
    return box


def _bullets_box(
    slide,
    bullets: list[str],
    *,
    left,
    top,
    width,
    height,
    template: PresentationTemplate,
):
    palette = template.palette
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    if not bullets:
        bullets = ["Materi akan dilengkapi."]

    for index, bullet in enumerate(bullets):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = 1.12
        paragraph.space_after = Pt(8)

        marker = paragraph.add_run()
        marker.text = "\u25aa  "  # small filled square
        marker.font.size = Pt(16)
        marker.font.bold = True
        marker.font.name = template.body_font
        marker.font.color.rgb = _rgb(palette.accent)

        run = paragraph.add_run()
        run.text = bullet
        run.font.size = Pt(18)
        run.font.name = template.body_font
        run.font.color.rgb = _rgb(palette.text)
    return box


# ── Slide builders ────────────────────────────────────────────────────────
def _build_cover_slide(
    prs: Presentation,
    template: PresentationTemplate,
    *,
    title: str,
    subtitle: str,
    audience: str,
    presenter: str,
    unit: str,
    date: str,
):
    palette = template.palette
    slide = _blank_slide(prs)
    style = template.cover_style
    logo_size = Inches(0.92)

    meta_lines = [line for line in (audience, presenter, unit, date) if line]
    meta_text = "  \u2022  ".join(meta_lines)

    if style == "split":
        _background(slide, palette.background)
        panel_w = Inches(5.0)
        _fill_rect(slide, palette.primary, 0, 0, panel_w, SLIDE_HEIGHT)
        add_brand_logo(slide, left=MARGIN, top=MARGIN, size=logo_size, fill_hex=palette.accent)
        _textbox(
            slide, BRAND_NAME, left=MARGIN, top=Inches(1.8), width=panel_w - MARGIN * 2,
            height=Inches(1.2), size=16, color=palette.on_primary, bold=True,
            font=template.heading_font, line_spacing=1.1,
        )
        if presenter or unit:
            _textbox(
                slide, "  \u2022  ".join([x for x in (presenter, unit) if x]),
                left=MARGIN, top=Inches(6.2), width=panel_w - MARGIN * 2, height=Inches(0.8),
                size=12, color=palette.on_primary, font=template.body_font,
            )
        _textbox(
            slide, title, left=panel_w + Inches(0.5), top=Inches(2.1),
            width=SLIDE_WIDTH - panel_w - Inches(1.0), height=Inches(2.4),
            size=40, color=palette.primary, bold=True, font=template.heading_font,
            line_spacing=1.05,
        )
        if subtitle:
            _textbox(
                slide, subtitle, left=panel_w + Inches(0.5), top=Inches(4.4),
                width=SLIDE_WIDTH - panel_w - Inches(1.0), height=Inches(1.0),
                size=18, color=palette.muted, font=template.body_font,
            )
        if meta_text:
            _textbox(
                slide, meta_text, left=panel_w + Inches(0.5), top=Inches(6.4),
                width=SLIDE_WIDTH - panel_w - Inches(1.0), height=Inches(0.6),
                size=12, color=palette.muted, font=template.body_font,
            )
        return

    if style == "sidebar":
        _background(slide, palette.background)
        bar_w = Inches(1.5)
        _fill_rect(slide, palette.primary, 0, 0, bar_w, SLIDE_HEIGHT)
        add_brand_logo(slide, left=Inches(0.29), top=MARGIN, size=logo_size, fill_hex=palette.accent)
        add_icon(slide, "cover", left=Inches(0.45), top=Inches(6.3), size=Inches(0.6), fill_hex=palette.accent)
        content_left = bar_w + Inches(0.6)
        _textbox(
            slide, BRAND_NAME, left=content_left, top=Inches(1.0), width=Inches(9.5),
            height=Inches(0.7), size=16, color=palette.primary, bold=True,
            font=template.heading_font,
        )
        _textbox(
            slide, title, left=content_left, top=Inches(2.2), width=Inches(10.0),
            height=Inches(2.2), size=40, color=palette.text, bold=True,
            font=template.heading_font, line_spacing=1.05,
        )
        _fill_rect(slide, palette.accent, content_left, Inches(4.45), Inches(2.2), Inches(0.08))
        if subtitle:
            _textbox(
                slide, subtitle, left=content_left, top=Inches(4.7), width=Inches(10.0),
                height=Inches(1.0), size=18, color=palette.muted, font=template.body_font,
            )
        if meta_text:
            _textbox(
                slide, meta_text, left=content_left, top=Inches(6.4), width=Inches(10.0),
                height=Inches(0.6), size=12, color=palette.muted, font=template.body_font,
            )
        return

    if style == "minimal":
        _background(slide, palette.background)
        add_brand_logo(slide, left=MARGIN, top=MARGIN, size=logo_size, fill_hex=palette.primary)
        _textbox(
            slide, BRAND_NAME, left=MARGIN + logo_size + Inches(0.2), top=MARGIN + Inches(0.18),
            width=Inches(8.0), height=Inches(0.6), size=15, color=palette.muted, bold=True,
            font=template.heading_font,
        )
        _textbox(
            slide, title, left=MARGIN, top=Inches(2.7), width=Inches(11.8), height=Inches(2.0),
            size=44, color=palette.text, bold=True, font=template.heading_font, line_spacing=1.04,
        )
        _fill_rect(slide, palette.accent, MARGIN, Inches(4.7), Inches(2.6), Inches(0.06))
        if subtitle:
            _textbox(
                slide, subtitle, left=MARGIN, top=Inches(4.95), width=Inches(11.8),
                height=Inches(1.0), size=18, color=palette.muted, font=template.body_font,
            )
        if meta_text:
            _textbox(
                slide, meta_text, left=MARGIN, top=Inches(6.5), width=Inches(11.8),
                height=Inches(0.6), size=12, color=palette.muted, font=template.body_font,
            )
        return

    if style == "frame":
        _background(slide, palette.background)
        # Bingkai tipis penuh + aksen.
        _fill_rect(slide, palette.primary, 0, 0, SLIDE_WIDTH, Inches(0.18))
        _fill_rect(slide, palette.primary, 0, SLIDE_HEIGHT - Inches(0.18), SLIDE_WIDTH, Inches(0.18))
        _fill_rect(slide, palette.accent, 0, Inches(0.18), Inches(0.18), SLIDE_HEIGHT - Inches(0.36))
        add_brand_logo(slide, left=MARGIN, top=MARGIN, size=logo_size, fill_hex=palette.primary)
        _textbox(
            slide, BRAND_NAME, left=MARGIN, top=Inches(1.7), width=Inches(11.8), height=Inches(0.7),
            size=16, color=palette.primary, bold=True, align=PP_ALIGN.CENTER, font=template.heading_font,
        )
        _textbox(
            slide, title, left=Inches(1.2), top=Inches(2.9), width=Inches(10.9), height=Inches(2.0),
            size=42, color=palette.text, bold=True, align=PP_ALIGN.CENTER,
            font=template.heading_font, line_spacing=1.05,
        )
        if subtitle:
            _textbox(
                slide, subtitle, left=Inches(1.2), top=Inches(4.9), width=Inches(10.9), height=Inches(1.0),
                size=18, color=palette.muted, align=PP_ALIGN.CENTER, font=template.body_font,
            )
        if meta_text:
            _textbox(
                slide, meta_text, left=Inches(1.2), top=Inches(6.4), width=Inches(10.9), height=Inches(0.6),
                size=12, color=palette.muted, align=PP_ALIGN.CENTER, font=template.body_font,
            )
        return

    # Default: 'band'
    _background(slide, palette.background)
    band_h = Inches(3.0)
    _fill_rect(slide, palette.primary, 0, 0, SLIDE_WIDTH, band_h)
    _fill_rect(slide, palette.accent, 0, band_h, SLIDE_WIDTH, Inches(0.14))
    add_brand_logo(slide, left=MARGIN, top=MARGIN, size=logo_size, fill_hex=palette.accent)
    _textbox(
        slide, BRAND_NAME, left=MARGIN + logo_size + Inches(0.2), top=MARGIN + Inches(0.2),
        width=Inches(9.0), height=Inches(0.6), size=16, color=palette.on_primary, bold=True,
        font=template.heading_font,
    )
    _textbox(
        slide, title, left=MARGIN, top=Inches(1.55), width=Inches(11.8), height=Inches(1.3),
        size=38, color=palette.on_primary, bold=True, font=template.heading_font, line_spacing=1.04,
    )
    if subtitle:
        _textbox(
            slide, subtitle, left=MARGIN, top=Inches(3.5), width=Inches(11.8), height=Inches(1.0),
            size=18, color=palette.text, font=template.body_font,
        )
    if meta_text:
        _textbox(
            slide, meta_text, left=MARGIN, top=Inches(6.5), width=Inches(11.8), height=Inches(0.6),
            size=12, color=palette.muted, font=template.body_font,
        )


def _add_header(slide, template: PresentationTemplate, header: str):
    palette = template.palette
    _fill_rect(slide, palette.surface, 0, 0, SLIDE_WIDTH, Inches(0.92))
    _fill_rect(slide, palette.primary, 0, Inches(0.92), SLIDE_WIDTH, Inches(0.04))
    add_brand_logo(slide, left=MARGIN, top=Inches(0.16), size=Inches(0.6), fill_hex=palette.primary)
    _textbox(
        slide, header, left=MARGIN + Inches(0.78), top=Inches(0.18), width=Inches(11.5),
        height=Inches(0.56), size=13, color=palette.muted, bold=True,
        font=template.body_font, anchor=MSO_ANCHOR.MIDDLE,
    )


def _add_footer(slide, template: PresentationTemplate, footer: str, page_number: int, total_pages: int):
    palette = template.palette
    _fill_rect(slide, palette.primary, 0, SLIDE_HEIGHT - Inches(0.42), SLIDE_WIDTH, Inches(0.42))
    _textbox(
        slide, footer, left=MARGIN, top=SLIDE_HEIGHT - Inches(0.42), width=Inches(10.6),
        height=Inches(0.42), size=10, color=palette.on_primary, font=template.body_font,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _textbox(
        slide, f"{page_number}/{total_pages}", left=SLIDE_WIDTH - Inches(1.4),
        top=SLIDE_HEIGHT - Inches(0.42), width=Inches(0.9), height=Inches(0.42), size=10,
        color=palette.on_primary, align=PP_ALIGN.RIGHT, font=template.body_font,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def _build_content_slide(
    prs: Presentation,
    template: PresentationTemplate,
    content: SlideContent,
    *,
    header: str,
    footer: str,
    page_number: int,
    total_pages: int,
):
    palette = template.palette
    slide = _blank_slide(prs)
    _background(slide, palette.background)
    _add_header(slide, template, header)

    icon_name = _CONTENT_ICON_CYCLE[(page_number - 1) % len(_CONTENT_ICON_CYCLE)]
    add_icon(slide, icon_name, left=MARGIN, top=Inches(1.35), size=Inches(0.42), fill_hex=palette.accent)

    _textbox(
        slide, content.title, left=MARGIN + Inches(0.6), top=Inches(1.2), width=Inches(11.5),
        height=Inches(0.9), size=26, color=palette.primary, bold=True, font=template.heading_font,
    )
    _fill_rect(slide, palette.accent, MARGIN, Inches(2.15), Inches(1.6), Inches(0.05))

    _bullets_box(
        slide, content.bullets, left=MARGIN, top=Inches(2.5), width=Inches(12.1),
        height=Inches(4.2), template=template,
    )

    _add_footer(slide, template, footer, page_number, total_pages)


def _build_closing_slide(
    prs: Presentation,
    template: PresentationTemplate,
    *,
    header: str,
    footer: str,
    presenter: str,
    unit: str,
    page_number: int,
    total_pages: int,
):
    palette = template.palette
    slide = _blank_slide(prs)
    _background(slide, palette.primary)
    add_brand_logo(slide, left=MARGIN, top=MARGIN, size=Inches(0.8), fill_hex=palette.accent)
    _textbox(
        slide, "Terima Kasih", left=Inches(1.0), top=Inches(2.7), width=Inches(11.3),
        height=Inches(1.4), size=46, color=palette.on_primary, bold=True, align=PP_ALIGN.CENTER,
        font=template.heading_font,
    )
    _textbox(
        slide, BRAND_NAME, left=Inches(1.0), top=Inches(4.2), width=Inches(11.3), height=Inches(0.7),
        size=18, color=palette.on_primary, align=PP_ALIGN.CENTER, font=template.body_font,
    )
    closer = "  \u2022  ".join([x for x in (presenter, unit) if x])
    if closer:
        _textbox(
            slide, closer, left=Inches(1.0), top=Inches(5.0), width=Inches(11.3), height=Inches(0.6),
            size=13, color=palette.accent, align=PP_ALIGN.CENTER, font=template.body_font,
        )
    _add_footer(slide, template, footer, page_number, total_pages)
