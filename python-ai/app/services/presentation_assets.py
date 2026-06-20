"""Local, no-internet asset registry for the presentation renderer (#221).

Semua aset visual (logo emblem + ikon) digambar sebagai shape vektor PPTX
secara deterministik. Tidak ada pengambilan gambar dari internet dan tidak ada
ketergantungan file biner eksternal, sehingga renderer aman dipakai di
lingkungan tanpa jaringan (baseline #225 / #227 menunda asset enrichment).
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# Identitas brand wajib (lihat #218/#221).
BRAND_NAME = "Istana Kepresidenan Yogyakarta"
BRAND_MONOGRAM = "IKY"

# Mode aset MVP (#225): seluruh aset visual digambar lokal (shape vektor),
# tidak ada pengambilan dari internet. Enrichment aset web ditunda ke #227.
ASSET_MODE = "local_assets_only"

# Nama shape logo dipakai sebagai marker agar bisa diverifikasi di test/QA.
LOGO_SHAPE_NAME = "ista-logo"

# Registry ikon lokal: nama logis -> auto shape PPTX. Dipakai sebagai aksen
# visual per jenis slide tanpa file gambar eksternal.
LOCAL_ICONS: dict[str, MSO_SHAPE] = {
    "cover": MSO_SHAPE.OVAL,
    "agenda": MSO_SHAPE.ROUNDED_RECTANGLE,
    "summary": MSO_SHAPE.RECTANGLE,
    "key_points": MSO_SHAPE.CHEVRON,
    "data": MSO_SHAPE.RECTANGLE,
    "activity": MSO_SHAPE.OVAL,
    "closing": MSO_SHAPE.OVAL,
}


def local_icon_names() -> tuple[str, ...]:
    """Daftar nama ikon lokal yang tersedia (untuk registry/QA no-internet)."""
    return tuple(LOCAL_ICONS.keys())


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.lstrip("#").upper())


def add_brand_logo(
    slide,
    *,
    left: Emu,
    top: Emu,
    size: Emu,
    fill_hex: str,
    text_hex: str = "FFFFFF",
) -> None:
    """Gambar emblem logo (oval + monogram) dan beri nama marker."""
    emblem = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    emblem.name = LOGO_SHAPE_NAME
    emblem.fill.solid()
    emblem.fill.fore_color.rgb = _rgb(fill_hex)
    emblem.line.color.rgb = _rgb(fill_hex)
    emblem.shadow.inherit = False

    text_frame = emblem.text_frame
    text_frame.word_wrap = False
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = BRAND_MONOGRAM
    run.font.bold = True
    run.font.size = Pt(max(8, int(size / Emu(Pt(1)) / 2.4)))
    run.font.color.rgb = _rgb(text_hex)
    run.font.name = "Arial"


def add_icon(
    slide,
    name: str,
    *,
    left: Emu,
    top: Emu,
    size: Emu,
    fill_hex: str,
) -> None:
    """Tempel ikon lokal sesuai registry. Nama tak dikenal -> fallback aksen."""
    shape_type = LOCAL_ICONS.get(name, MSO_SHAPE.RECTANGLE)
    icon = slide.shapes.add_shape(shape_type, left, top, size, size)
    icon.name = f"ista-icon-{name}"
    icon.fill.solid()
    icon.fill.fore_color.rgb = _rgb(fill_hex)
    icon.line.fill.background()
    icon.shadow.inherit = False
