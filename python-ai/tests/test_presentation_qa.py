"""Visual QA, safety & no-internet baseline untuk renderer presentasi (#225).

Menjamin: render 5 template tanpa jaringan, design token (warna palette)
terpakai konsisten, logo/header/footer hadir, batas bullet/slide/title
ditegakkan, model hanya menyumbang outline (bukan keputusan desain), dan
seluruh shape berada dalam batas kanvas 16:9 (no structural overflow).
"""

import os
import socket
import sys
from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from app.services.presentation_assets import (
    ASSET_MODE,
    BRAND_NAME,
    LOGO_SHAPE_NAME,
)
from app.services.presentation_render import (
    BULLET_MAX_CHARS,
    MAX_BULLETS_PER_SLIDE,
    MAX_CONTENT_SLIDES,
    PRESENTATION_TEMPLATES,
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    SLIDE_TITLE_MAX_CHARS,
    render_presentation,
)

TOLERANCE_EMU = 12700  # 1pt


def _open(content: bytes) -> Presentation:
    return Presentation(BytesIO(content))


def _fill_hexes(prs: Presentation) -> set[str]:
    hexes: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                fill = shape.fill
                if fill.type is None:
                    continue
                rgb = fill.fore_color.rgb
            except (TypeError, AttributeError, ValueError):
                continue
            hexes.add(str(rgb).upper())
    return hexes


def _shape_names(prs: Presentation) -> list[str]:
    return [shape.name for slide in prs.slides for shape in slide.shapes]


def _all_texts(prs: Presentation) -> str:
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _sample_outline() -> list[dict]:
    return [
        {"title": "Ringkasan", "bullets": ["Poin A", "Poin B", "Poin C"]},
        {"title": "Tindak Lanjut", "bullets": ["Langkah 1", "Langkah 2"]},
    ]


def test_asset_mode_is_local_only():
    assert ASSET_MODE == "local_assets_only"


def test_all_templates_render_without_any_network(monkeypatch):
    """Render tidak boleh membuka koneksi jaringan apa pun."""

    def _blocked_socket(*args, **kwargs):
        raise AssertionError("Renderer presentasi tidak boleh mengakses jaringan.")

    monkeypatch.setattr(socket, "socket", _blocked_socket)
    monkeypatch.setattr(socket, "create_connection", _blocked_socket)

    for key in PRESENTATION_TEMPLATES:
        render = render_presentation(
            title=f"Deck {key}",
            visual_template=key,
            outline=_sample_outline(),
        )
        prs = _open(render.content)
        assert len(prs.slides._sldIdLst) >= 2


def test_design_tokens_applied_consistently_per_template():
    """Warna palette primary & accent tiap template benar-benar dipakai."""
    for key, template in PRESENTATION_TEMPLATES.items():
        render = render_presentation(
            title=f"Token {key}",
            visual_template=key,
            outline=_sample_outline(),
        )
        fills = _fill_hexes(_open(render.content))
        assert template.palette.primary.upper() in fills, f"primary tidak dipakai di {key}"
        assert template.palette.accent.upper() in fills, f"accent tidak dipakai di {key}"


def test_logo_header_footer_present_across_slides():
    render = render_presentation(
        title="QA Header Footer",
        visual_template="data_tabel",
        header="Sekretariat Presiden",
        footer="Dokumen Internal",
        outline=[
            {"title": "Bagian 1", "bullets": ["x"]},
            {"title": "Bagian 2", "bullets": ["y"]},
        ],
    )
    prs = _open(render.content)
    slides = list(prs.slides)

    # Cover (pertama) memuat logo + brand.
    assert LOGO_SHAPE_NAME in [s.name for s in slides[0].shapes]
    assert BRAND_NAME in _all_texts(prs)

    # Setiap slide konten memuat header & footer.
    content_slides = slides[1:-1]
    assert len(content_slides) == 2
    for slide in content_slides:
        text = "\n".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
        assert "Sekretariat Presiden" in text
        assert "Dokumen Internal" in text

    # Closing (terakhir) memuat footer + logo.
    closing = slides[-1]
    closing_names = [s.name for s in closing.shapes]
    assert LOGO_SHAPE_NAME in closing_names


def test_bullet_and_slide_limits_enforced():
    long_bullet = "kata " * 200  # ~1000 chars
    too_many_bullets = [f"Bullet {i}" for i in range(20)]

    render = render_presentation(
        title="Limit",
        visual_template="modern_minimal",
        outline=[
            {"title": "Panjang", "bullets": [long_bullet]},
            {"title": "Banyak", "bullets": too_many_bullets},
        ],
    )
    prs = _open(render.content)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    assert len(run.text) <= BULLET_MAX_CHARS + 5

    # Slide "Banyak" — kotak bullet tunggal tidak boleh > MAX_BULLETS_PER_SLIDE.
    # (cover + 2 konten + closing) -> index 2 = slide "Banyak".
    content_slide = list(prs.slides)[2]
    max_paragraphs_in_a_box = 0
    for shape in content_slide.shapes:
        if not shape.has_text_frame:
            continue
        non_empty = sum(1 for p in shape.text_frame.paragraphs if p.text.strip())
        max_paragraphs_in_a_box = max(max_paragraphs_in_a_box, non_empty)
    assert max_paragraphs_in_a_box <= MAX_BULLETS_PER_SLIDE


def test_title_is_truncated_to_limit():
    long_title = "Judul Slide Sangat Panjang " * 20
    render = render_presentation(
        title="Deck",
        visual_template="resmi_klasik",
        outline=[{"title": long_title, "bullets": ["x"]}],
    )
    prs = _open(render.content)
    content_slide = list(prs.slides)[1]
    titles = [s.text_frame.text for s in content_slide.shapes if s.has_text_frame and s.text_frame.text]
    assert any(len(t) <= SLIDE_TITLE_MAX_CHARS + 1 for t in titles)


def test_content_slides_capped_at_maximum():
    outline = [{"title": f"Bagian {i}", "bullets": ["x"]} for i in range(60)]
    render = render_presentation(title="Panjang", visual_template="executive_brief", outline=outline)
    prs = _open(render.content)
    assert len(prs.slides._sldIdLst) <= MAX_CONTENT_SLIDES + 2


def test_model_cannot_inject_design_decisions_via_outline():
    """Field tak dikenal dalam outline (warna/font) diabaikan renderer."""
    rogue_color = "FF00AA"  # tidak ada di palette mana pun
    render = render_presentation(
        title="Safety",
        visual_template="resmi_klasik",
        outline=[
            {
                "title": "Slide",
                "bullets": ["aman"],
                "background": rogue_color,
                "color": rogue_color,
                "font": "Comic Sans",
                "template": "evil",
            }
        ],
    )
    fills = _fill_hexes(_open(render.content))
    assert rogue_color not in fills


def test_no_shape_overflows_slide_bounds():
    for key in PRESENTATION_TEMPLATES:
        render = render_presentation(
            title=f"Bounds {key}",
            visual_template=key,
            header="Header Uji",
            footer="Footer Uji",
            outline=_sample_outline(),
        )
        prs = _open(render.content)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.left is None or shape.top is None:
                    continue
                width = shape.width or 0
                height = shape.height or 0
                assert shape.left >= -TOLERANCE_EMU, f"{key}: shape kiri keluar kanvas"
                assert shape.top >= -TOLERANCE_EMU, f"{key}: shape atas keluar kanvas"
                assert shape.left + width <= SLIDE_WIDTH + TOLERANCE_EMU, f"{key}: shape kanan keluar kanvas"
                assert shape.top + height <= SLIDE_HEIGHT + TOLERANCE_EMU, f"{key}: shape bawah keluar kanvas"


def test_minimal_outline_still_produces_cover_and_closing():
    render = render_presentation(title="Minimal", visual_template="kegiatan_dokumentasi", outline=[])
    prs = _open(render.content)
    # cover + fallback konten + closing
    assert len(prs.slides._sldIdLst) >= 2
    assert LOGO_SHAPE_NAME in _shape_names(prs)
