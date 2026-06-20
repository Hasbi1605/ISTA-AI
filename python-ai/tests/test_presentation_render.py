import os
import sys
from io import BytesIO
from pathlib import Path

from fastapi.testclient import TestClient
from pptx import Presentation

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
os.environ["AI_SERVICE_TOKEN"] = "test_internal_api_secret"

from app.documents_api import app
from app.services.presentation_assets import BRAND_NAME, LOGO_SHAPE_NAME, local_icon_names
from app.services.presentation_render import (
    BULLET_MAX_CHARS,
    MAX_CONTENT_SLIDES,
    PRESENTATION_TEMPLATES,
    render_presentation,
)

AUTH = {"Authorization": "Bearer test_internal_api_secret"}


def _client() -> TestClient:
    return TestClient(app)


def _open(content: bytes) -> Presentation:
    return Presentation(BytesIO(content))


def _all_texts(prs: Presentation) -> str:
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _shape_names(prs: Presentation) -> list[str]:
    names = []
    for slide in prs.slides:
        for shape in slide.shapes:
            names.append(shape.name)
    return names


# ── Routing & dependency ────────────────────────────────────────────────────
def test_presentation_generate_route_is_registered_on_document_app():
    paths = {route.path for route in app.routes}
    assert "/api/presentations/generate" in paths
    assert "/api/presentations/templates" in paths


def test_presentation_route_absent_from_chat_app():
    from app.chat_api import app as chat_app

    paths = {route.path for route in chat_app.routes}
    assert "/api/presentations/generate" not in paths


def test_python_pptx_is_listed_in_requirements():
    requirements = (Path(__file__).resolve().parents[1] / "requirements.txt").read_text(
        encoding="utf-8"
    )
    normalized = {
        line.strip().lower().split("==", 1)[0]
        for line in requirements.splitlines()
        if line.strip() and not line.lstrip().startswith("#")
    }
    assert "python-pptx" in normalized


# ── Renderer service ────────────────────────────────────────────────────────
def test_render_minimal_outline_produces_valid_pptx():
    render = render_presentation(
        title="Paparan Kesiapan Kegiatan",
        visual_template="modern_minimal",
        outline=[{"title": "Agenda", "bullets": ["Pembukaan", "Pembahasan", "Penutup"]}],
    )

    assert render.filename.endswith(".pptx")
    assert render.template == "modern_minimal"

    prs = _open(render.content)
    # cover + 1 konten + closing
    assert len(prs.slides._sldIdLst) == 3
    assert render.slide_count == 3


def test_logo_header_footer_present_on_render():
    render = render_presentation(
        title="Rapat Koordinasi Internal",
        visual_template="resmi_klasik",
        header="Sekretariat Presiden - Yogyakarta",
        footer="Dokumen Internal ISTA AI",
        presenter="Bagian Umum",
        outline=[{"title": "Latar Belakang", "bullets": ["Poin satu", "Poin dua"]}],
    )

    prs = _open(render.content)
    names = _shape_names(prs)
    texts = _all_texts(prs)

    assert LOGO_SHAPE_NAME in names
    assert BRAND_NAME in texts
    assert "Sekretariat Presiden - Yogyakarta" in texts
    assert "Dokumen Internal ISTA AI" in texts


def test_all_five_templates_render_without_internet():
    assert set(PRESENTATION_TEMPLATES.keys()) == {
        "resmi_klasik",
        "modern_minimal",
        "executive_brief",
        "data_tabel",
        "kegiatan_dokumentasi",
    }

    for key in PRESENTATION_TEMPLATES:
        render = render_presentation(
            title=f"Deck {key}",
            visual_template=key,
            outline=[
                {"title": "Ringkasan", "bullets": ["A", "B"]},
                {"title": "Tindak Lanjut", "bullets": ["C"]},
            ],
        )
        prs = _open(render.content)
        assert len(prs.slides._sldIdLst) >= 2
        assert LOGO_SHAPE_NAME in _shape_names(prs)


def test_long_bullet_is_truncated_and_does_not_error():
    long_bullet = "Lorem ipsum dolor sit amet " * 80  # ~2000 chars
    render = render_presentation(
        title="Uji Bullet Panjang",
        visual_template="executive_brief",
        outline=[{"title": "Detail", "bullets": [long_bullet]}],
    )

    prs = _open(render.content)
    texts = _all_texts(prs)
    # Bullet harus terpangkas (ellipsis) dan tidak melebihi batas + marker.
    assert "\u2026" in texts
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        assert len(run.text) <= BULLET_MAX_CHARS + 5


def test_too_many_outline_slides_are_capped():
    outline = [{"title": f"Bagian {i}", "bullets": ["x"]} for i in range(40)]
    render = render_presentation(
        title="Deck Panjang",
        visual_template="data_tabel",
        outline=outline,
    )
    prs = _open(render.content)
    # cover + (<= MAX_CONTENT_SLIDES) + closing
    assert len(prs.slides._sldIdLst) <= MAX_CONTENT_SLIDES + 2


def test_slide_count_target_caps_content_slides():
    outline = [{"title": f"Bagian {i}", "bullets": ["x"]} for i in range(10)]
    render = render_presentation(
        title="Deck Ringkas",
        visual_template="modern_minimal",
        outline=outline,
        slide_count=4,
    )
    prs = _open(render.content)
    # cover + 2 konten + closing == 4
    assert len(prs.slides._sldIdLst) == 4


def test_unknown_template_is_rejected_by_service():
    try:
        render_presentation(title="X", visual_template="tidak_ada", outline=[])
    except ValueError as exc:
        assert "Template visual tidak dikenal" in str(exc)
    else:
        raise AssertionError("ValueError tidak dilempar untuk template tak dikenal")


def test_blank_title_is_rejected_by_service():
    try:
        render_presentation(title="   ", visual_template="resmi_klasik", outline=[])
    except ValueError as exc:
        assert "Judul presentasi wajib diisi" in str(exc)
    else:
        raise AssertionError("ValueError tidak dilempar untuk judul kosong")


def test_local_icon_registry_is_offline_and_named():
    names = local_icon_names()
    assert "cover" in names
    assert "agenda" in names
    assert len(names) >= 5


# ── Endpoint ──────────────────────────────────────────────────────────────
def test_generate_endpoint_returns_pptx_with_security_headers():
    response = _client().post(
        "/api/presentations/generate",
        headers=AUTH,
        json={
            "title": "Paparan Evaluasi Triwulan",
            "visual_template": "kegiatan_dokumentasi",
            "header": "Istana Kepresidenan Yogyakarta",
            "footer": "Rahasia Internal",
            "outline": [
                {"title": "Capaian", "bullets": ["Target tercapai", "Catatan"]},
                {"title": "Rencana", "bullets": ["Langkah berikutnya"]},
            ],
        },
    )

    assert response.status_code == 200
    assert (
        response.headers["content-type"]
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert response.headers["x-content-type-options"] == "nosniff"
    assert response.headers["cache-control"] == "no-store"
    assert response.headers["x-presentation-template"] == "kegiatan_dokumentasi"
    assert "attachment;" in response.headers["content-disposition"]

    prs = _open(response.content)
    assert len(prs.slides._sldIdLst) == 4  # cover + 2 + closing


def test_generate_endpoint_rejects_unknown_template_with_400():
    response = _client().post(
        "/api/presentations/generate",
        headers=AUTH,
        json={"title": "Deck", "visual_template": "alien_theme", "outline": []},
    )
    assert response.status_code == 400
    assert "Template visual tidak dikenal" in response.json()["detail"]


def test_generate_endpoint_rejects_blank_title_with_400():
    response = _client().post(
        "/api/presentations/generate",
        headers=AUTH,
        json={"title": "   ", "visual_template": "resmi_klasik", "outline": []},
    )
    assert response.status_code == 400


def test_generate_endpoint_requires_auth():
    response = _client().post(
        "/api/presentations/generate",
        json={"title": "Deck", "visual_template": "resmi_klasik", "outline": []},
    )
    assert response.status_code == 401
